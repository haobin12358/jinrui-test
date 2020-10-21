import oss2, os, cv2
from pptx import Presentation, slide
from pptx.util import Inches, Pt
from jinrui.config.secret import ACCESS_KEY_ID, ACCESS_KEY_SECRET, ALIOSS_BUCKET_NAME, ALIOSS_ENDPOINT

def get_number(char):
    """
    判断字符串中，中文的个数
    :param char: 字符串
    :return:
    """
    count = 0
    for item in char:
        if 0x4E00 <= ord(item) <= 0x9FA5:
            count += 1
    return count

auth = oss2.Auth(ACCESS_KEY_ID, ACCESS_KEY_SECRET)
bucket = oss2.Bucket(auth, ALIOSS_ENDPOINT, ALIOSS_BUCKET_NAME)


if __name__ == "__main__":
    question = "<div>9．【图形变换】如图6，Rt△ABC中，∠BAC＝90°，AB＝AC，将△ABC绕点C顺时针旋转40°得到△A′B′C，CB′与AB相交于点D，连结AA′，则∠B′A′A的度数为()</div><div>A．10° B．15° C．20° D．30°</div>"
    answer = "<div>1．C【解析】 ∵π＞<img src='https://jinrui-sheet.oss-cn-shanghai.aliyuncs.com/image1_5881953458786304.png'></img>＞0＞－6，∴所给的各数中最大的数是π.</div>"
    knowledge = "实数的有关概念与大小比较"

    question_dict = question.split("<div>")

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    # 创建背景
    img_path = os.path.abspath("control/ppt.jpg")
    back_pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    slide.shapes._spTree.remove(back_pic._element)
    slide.shapes._spTree.insert(2, back_pic._element)

    left = top = Inches(0.5)
    for question_item in question_dict:
        print(">>>>>>>>>>>>>>>>>>question_item:" + question_item)
        if question_item:
            question_item = question_item.replace("</div>", "#####").replace("<img src='", "#####").replace("'></img>", "#####")
            question_item_dict = question_item.split("#####")
            width = Inches(1)
            height = Inches(0.4)
            for row in question_item_dict:
                if row:
                    if "https://" in row:
                        print(len(row))
                        row_dict = row.split("/")
                        pic_save_path = "D:\\" + row_dict[-1]
                        bucket.get_object_to_file(row_dict[-1], pic_save_path)
                        img = cv2.imread(pic_save_path)
                        pic_width = img.shape[0]
                        pic_height = img.shape[1]
                        pic = slide.shapes.add_picture(pic_save_path, left, top, height=height)
                        left = left + pic_width * height / pic_height
                        os.remove(pic_save_path)
                    else:
                        width = Inches(len(row) / 4 + 2)
                        print(width)
                        print(Inches(12))
                        if left + width > Inches(12):
                            use_width = Inches(12) - left
                            print(use_width)
                            from pptx.util import Length
                            use_word = Length(4 * use_width).inches
                            print(">>>>>>>>>>>>>>>>>use_word:" + str(use_word))
                            txBox = slide.shapes.add_textbox(left, top, use_width, height)
                            tf = txBox.text_frame
                            print(row)
                            tf.text = row[0: int(use_word)]
                            left = Inches(0.5)
                            top = top + Inches(0.5)
                            other_word_num = len(row) - use_word
                            if other_word_num % 36 == 0:
                                perform = other_word_num / 36
                            else:
                                perform = int(other_word_num / 36) + 1
                            i = 0
                            while i < perform:
                                txBox = slide.shapes.add_textbox(left, top,
                                                                 Inches(len(row[int(use_word) + 1 + 36 * i:
                                                                                int(use_word) + 36 * (i + 1) + 1]) / 4 + 2),
                                                                 height)
                                tf = txBox.text_frame
                                tf.text = row[int(use_word) + 1 + 36 * i: int(use_word) + 36 * (i + 1) + 1]
                                left = Inches(0.5)
                                top = top + Inches(0.5)
                                i += 1
                        else:
                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            tf = txBox.text_frame
                            tf.text = row
                            left = left + width
            left = Inches(0.5)
            top = top + Inches(0.5)
    prs.save('D:\\test.pptx')

    e = Inches(12)