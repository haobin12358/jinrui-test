import oss2, os, cv2, platform, uuid, shutil

from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from flask import current_app

from ..extensions.success_response import Success
from ..extensions.params_validates import parameter_required
from jinrui.extensions.register_ext import db
from jinrui.extensions.register_ext import ali_oss
from jinrui.models.jinrui import j_question, j_paper, j_student, j_answer_booklet, j_score
from jinrui.config.secret import ACCESS_KEY_ID, ACCESS_KEY_SECRET, ALIOSS_BUCKET_NAME, ALIOSS_ENDPOINT

class CPpt():

    def get_wrong_paper_ppt(self):
        """
        获取错题ppt
        """
        current_app.logger.info(datetime.now())
        args = parameter_required(("paperName", "classId", "errRate"))

        # 获取试卷
        paper = j_paper.query.filter(j_paper.name == args.get("paperName")).first_("未找到该试卷")
        paper_id = paper.id

        # 获取题目列表
        question_list = j_question.query.filter(j_question.paper_id == paper_id).all()

        # 题号list
        question_num_list = []
        # 分数list
        question_score_list = []
        for question in question_list:
            question_num_list.append(question.question_number)
            question_score_list.append(question.score)

        current_app.logger.info(">>>>>>>>>>>>>>>>>>>>>>>>>question_num_list:" + str(question_num_list))
        current_app.logger.info(">>>>>>>>>>>>>>>>>>>>>>>>>question_score_list:" + str(question_score_list))
        # 学生list
        student_id_list = []
        student_list = j_student.query.filter(j_student.org_id == args.get("classId")).all()
        for student in student_list:
            student_id_list.append(student.id)

        current_app.logger.info(">>>>>>>>>>>>>>>>>>>>>>>>>student_id_list:" + str(student_id_list))
        answer_id_list = []
        for student_id in student_id_list:
            filter_args = [j_answer_booklet.paper_id == paper_id, j_answer_booklet.student_id == student_id]
            if args.get("startTime"):
                filter_args.append(j_answer_booklet.create_time >= args.get("startTime"))
            if args.get("endTime"):
                filter_args.append(j_answer_booklet.create_time <= args.get("endTime"))
            answer_booklet_list = j_answer_booklet.query.filter(*filter_args).all()
            for answer_booklet_item in answer_booklet_list:
                answer_id_list.append(answer_booklet_item.id)

        current_app.logger.info(">>>>>>>>>>>>>>>>>>>>>>>>>answer_id_list:" + str(answer_id_list))
        # 总答卷数目
        total_answer_booklet = len(answer_id_list)
        # 要生成ppt的question列表
        ppt_question_number_list = []
        i = 0
        while i < len(question_num_list):
            j = 0
            for answer_id in answer_id_list:
                score = j_score.query.filter(j_score.booklet_id == answer_id,
                                             j_score.question_number == int(question_num_list[i]),
                                             j_score.score == question_score_list[i])\
                    .first()
                if score:
                    j = j + 1
            if (total_answer_booklet - j) / total_answer_booklet > float(args.get("errRate")) / 100:
                ppt_question_number_list.append(int(question_num_list[i]))
            i = i + 1

        ppt_question_number_list.sort()

        # ppt用到的题目答案列表
        question_answer_list = []
        for question_number in ppt_question_number_list:
            question_answer_dict = {}
            ppt_question_answer = j_question.query.filter(j_question.question_number == str(question_number),
                                                          j_question.paper_id == paper_id).first_("未找到题目")
            question_answer_dict["question"] = ppt_question_answer.content
            question_answer_dict["answer"] = ppt_question_answer.answer
            question_answer_dict["knowledge"] = ppt_question_answer.knowledge

            question_answer_list.append(question_answer_dict)

        current_app.logger.info(">>>>>>>>>>>>>>>>>>>question_answer_list:" + str(question_answer_list))
        # 阿里云oss参数
        auth = oss2.Auth(ACCESS_KEY_ID, ACCESS_KEY_SECRET)
        bucket = oss2.Bucket(auth, ALIOSS_ENDPOINT, ALIOSS_BUCKET_NAME)

        # 创建ppt
        prs = Presentation()
        # ppt样式
        blank_slide_layout = prs.slide_layouts[6]

        ppt_uuid = str(uuid.uuid1())
        if platform.system() == "Windows":
            pic_path = "D:\\jinrui_pic\\" + ppt_uuid + "\\"
        else:
            pic_path = "/tmp/jinrui_pic/" + ppt_uuid + "/"
        if not os.path.exists(pic_path):
            os.makedirs(pic_path)


        for question_answer in question_answer_list:
            # 创建幻灯片
            slide = prs.slides.add_slide(blank_slide_layout)

            question_dict = question_answer["question"].split("<div>")
            # 设置顶点尺寸
            left = top = Inches(0.5)


            for question_item in question_dict:
                if question_item:
                    # 创建背景
                    current_app.logger.info(os.path.abspath("jinrui/control/ppt.jpg"))
                    img_path = os.path.abspath("jinrui/control/ppt.jpg")
                    back_pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width,
                                                        height=prs.slide_height)
                    slide.shapes._spTree.remove(back_pic._element)
                    slide.shapes._spTree.insert(2, back_pic._element)

                    question_item = question_item.replace("</div>", "#####").replace("<img src='", "#####").replace(
                        "'></img>", "#####")
                    question_item_dict = question_item.split("#####")
                    width = Inches(1)
                    height = Inches(0.4)
                    for row in question_item_dict:
                        if row:
                            if "https://" in row:
                                row_dict = row.split("/")
                                pic_save_path = pic_path + row_dict[-1]
                                # 存储图片到本地
                                bucket.get_object_to_file(row_dict[-1], pic_save_path)
                                img = cv2.imread(pic_save_path)
                                pic_width = img.shape[0]
                                pic_height = img.shape[1]
                                # 添加图片到ppt中
                                pic = slide.shapes.add_picture(pic_save_path, left, top, height=height)
                                width = width + pic_width * height / pic_height
                                # 移除图片
                                os.remove(pic_save_path)
                            else:
                                # chinese_num = self._get_number(row)
                                # print(">>>>>>>>>>>>>>>>>>chinese_num:" + str(chinese_num))
                                # english_num = len(row) - chinese_num
                                # print(">>>>>>>>>>>>>>>>>>english_num:" + str(english_num))

                                width = Inches(len(row) / 4 + 2)
                                current_app.logger.info(">>>>>>>>>>>>>>>>>>width:" + str(width))
                                txBox = slide.shapes.add_textbox(left, top, width, height)
                                tf = txBox.text_frame
                                tf.text = row
                                left = left + width
                left = Inches(0.5)
                top = top + Inches(0.5)

            slide = prs.slides.add_slide(blank_slide_layout)
            answer_dict = question_answer["answer"].split("<div>")
            # 设置顶点尺寸
            left = top = Inches(0.5)

            for answer_item in answer_dict:
                if answer_item:
                    # 创建背景
                    current_app.logger.info(os.path.abspath("jinrui/control/ppt.jpg"))
                    img_path = os.path.abspath("jinrui/control/ppt.jpg")
                    back_pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width,
                                                        height=prs.slide_height)
                    slide.shapes._spTree.remove(back_pic._element)
                    slide.shapes._spTree.insert(2, back_pic._element)

                    answer_item = answer_item.replace("</div>", "#####").replace("<img src='", "#####").replace(
                        "'></img>", "#####")
                    question_item_dict = answer_item.split("#####")
                    width = Inches(1)
                    height = Inches(0.4)
                    for row in question_item_dict:
                        if row:
                            if "https://" in row:
                                row_dict = row.split("/")
                                pic_save_path = pic_path + row_dict[-1]
                                bucket.get_object_to_file(row_dict[-1], pic_save_path)
                                img = cv2.imread(pic_save_path)
                                pic_width = img.shape[0]
                                pic_height = img.shape[1]
                                pic = slide.shapes.add_picture(pic_save_path, left, top, height=height)
                                width = width + pic_width * height / pic_height
                                os.remove(pic_save_path)
                            else:
                                # chinese_num = self._get_number(row)
                                # print(">>>>>>>>>>>>>>>>>>chinese_num:" + str(chinese_num))
                                # english_num = len(row) - chinese_num
                                # print(">>>>>>>>>>>>>>>>>>english_num:" + str(english_num))

                                width = Inches(len(row) / 4 + 2)
                                current_app.logger.info(">>>>>>>>>>>>>>>>>>width:" + str(width))
                                txBox = slide.shapes.add_textbox(left, top, width, height)
                                tf = txBox.text_frame
                                tf.text = row
                                left = left + width
                left = Inches(0.5)
                top = top + Inches(0.5)
        prs.save(pic_path + "ppt-" + ppt_uuid + ".pptx")

        ppt_url = "https://" + ALIOSS_BUCKET_NAME + "." + ALIOSS_ENDPOINT + "/" + "ppt-" + ppt_uuid + ".pptx"
        result = bucket.put_object_from_file("ppt-" + ppt_uuid + ".pptx", pic_path + "ppt-" + ppt_uuid + ".pptx")
        current_app.logger.info(result.status)
        current_app.logger.info(ppt_url)
        shutil.rmtree(pic_path)
        current_app.logger.info(datetime.now())
        return Success(data={"url": ppt_url})

    def _get_number(self, div_str):
        """
        判断字符串中，中文的个数
        :param char: 字符串
        :return:
        """
        count = 0
        for item in div_str:
            if 0x4E00 <= ord(item) <= 0x9FA5:
                count += 1
        return count